#!/usr/bin/env python3
"""
PowerPoint Presentation Generator for Piano Keys Production Readiness Plan
Creates a beautiful, professional presentation with charts and visual design.
"""

from pathlib import Path
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re


def create_powerpoint_presentation(source_file: Path, output_dir: Path) -> Path:
    """
    Generate a beautiful PowerPoint presentation from the production readiness plan.

    Args:
        source_file: Path to the source markdown file
        output_dir: Directory to save the generated presentation

    Returns:
        Path to the generated .pptx file
    """
    # Create presentation
    prs = Presentation()

    # Set slide size (16:9 widescreen)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Read source content
    with open(source_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Parse content into sections
    sections = parse_markdown_sections(content)

    # Add slides
    add_title_slide(prs)
    add_agenda_slide(prs, sections)
    add_executive_summary_slides(prs, sections)
    add_phase_slides(prs, sections)
    add_closing_slide(prs)

    # Save presentation
    output_path = output_dir / "production_readiness_presentation.pptx"
    prs.save(output_path)

    return output_path


def parse_markdown_sections(content: str) -> dict:
    """Parse markdown content into structured sections."""
    sections = {
        'title': 'Piano Keys Production Readiness Plan',
        'subtitle': '12-Week Roadmap to Launch',
        'executive_summary': {},
        'phases': [],
        'metrics': {},
        'resources': {}
    }

    lines = content.split('\n')

    # Extract key information
    for i, line in enumerate(lines):
        if '**Current State**' in line:
            sections['executive_summary']['current'] = '\n'.join(lines[i:i+6])
        elif '**Target State**' in line:
            sections['executive_summary']['target'] = '\n'.join(lines[i:i+7])
        elif line.startswith('### Phase'):
            phase_title = line.lstrip('#').strip()
            # Get next few lines for phase description
            phase_desc = []
            j = i + 1
            while j < len(lines) and j < i + 10 and not lines[j].startswith('#'):
                if lines[j].strip():
                    phase_desc.append(lines[j].strip())
                j += 1
            sections['phases'].append({
                'title': phase_title,
                'description': '\n'.join(phase_desc)
            })

    return sections


def add_title_slide(prs):
    """Add professional title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Piano Keys"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(31, 71, 136)  # #1F4788
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Production Readiness Roadmap"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(46, 92, 138)  # #2E5C8A
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Add date
    date_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.5), Inches(8), Inches(0.5)
    )
    date_frame = date_box.text_frame
    date_frame.text = date.today().strftime('%B %d, %Y')
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(16)
    date_para.alignment = PP_ALIGN.CENTER


def add_agenda_slide(prs, sections):
    """Add agenda/table of contents slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content

    # Set title
    title = slide.shapes.title
    title.text = "Agenda"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 71, 136)

    # Add content
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    agenda_items = [
        "Executive Summary",
        "6-Phase Implementation Roadmap",
        "Critical Security Issues",
        "Resource Requirements",
        "Success Metrics",
        "Next Steps"
    ]

    for item in agenda_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(24)


def add_executive_summary_slides(prs, sections):
    """Add executive summary slides."""
    # Current State slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Current State: 30% Production Ready"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 71, 136)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    current_items = [
        "CRITICAL: Authentication bypassed",
        "CRITICAL: Hardcoded secret keys",
        "Testing coverage: 12-15% (need 80%)",
        "SQLite database (production blocker)",
        "No CI/CD pipeline",
        "Minimal accessibility"
    ]

    for item in current_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        if "CRITICAL" in item:
            p.font.color.rgb = RGBColor(220, 53, 69)  # Red
            p.font.bold = True

    # Target State slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Target State: 100% Production Ready"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 169, 224)  # Cyan

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    target_items = [
        "✓ Zero security vulnerabilities",
        "✓ 80%+ test coverage",
        "✓ PostgreSQL + S3/MinIO + Redis",
        "✓ Full CI/CD pipeline",
        "✓ WCAG AA accessibility",
        "✓ Comprehensive monitoring"
    ]

    for item in target_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(40, 167, 69)  # Green
        p.font.bold = True


def add_phase_slides(prs, sections):
    """Add slides for each phase."""
    phases = [
        {
            'num': 1,
            'title': 'Security & Foundation',
            'weeks': 'Weeks 1-2',
            'icon': '🔒',
            'goal': 'Eliminate critical security vulnerabilities',
            'deliverables': [
                'Fix authentication bypass',
                'Replace hardcoded secrets',
                'Add rate limiting',
                'Implement security headers',
                'Structured logging'
            ]
        },
        {
            'num': 2,
            'title': 'Infrastructure Migration',
            'weeks': 'Weeks 3-4',
            'icon': '🏗️',
            'goal': 'Migrate to production-grade database and storage',
            'deliverables': [
                'PostgreSQL migration',
                'Redis caching layer',
                'S3/MinIO file storage',
                'Horizontal scaling',
                'Database optimization'
            ]
        },
        {
            'num': 3,
            'title': 'Testing & Quality',
            'weeks': 'Weeks 5-7',
            'icon': '✅',
            'goal': 'Achieve 80% test coverage',
            'deliverables': [
                '215+ backend tests',
                '150+ frontend tests',
                '20+ E2E tests',
                'Code quality tools',
                'Pre-commit hooks'
            ]
        },
        {
            'num': 4,
            'title': 'User Experience',
            'weeks': 'Weeks 8-9',
            'icon': '👥',
            'goal': 'User-centered design with WCAG AA compliance',
            'deliverables': [
                'User authentication UI',
                'Onboarding wizard',
                'Accessibility audit',
                'Notification system',
                'Help system'
            ]
        },
        {
            'num': 5,
            'title': 'CI/CD & Monitoring',
            'weeks': 'Weeks 10-11',
            'icon': '📊',
            'goal': 'Automated deployment and full observability',
            'deliverables': [
                'GitHub Actions CI/CD',
                'Prometheus metrics',
                'Grafana dashboards',
                'Sentry error tracking',
                'Operational runbooks'
            ]
        },
        {
            'num': 6,
            'title': 'Production Hardening',
            'weeks': 'Week 12',
            'icon': '🚀',
            'goal': 'Final testing and launch preparation',
            'deliverables': [
                'Security audit (OWASP)',
                'Load testing (100+ users)',
                'Disaster recovery testing',
                'Launch checklist',
                'GO LIVE!'
            ]
        }
    ]

    for phase in phases:
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Title with icon
        title = slide.shapes.title
        title.text = f"{phase['icon']} Phase {phase['num']}: {phase['title']}"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 71, 136)

        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()

        # Weeks
        p = tf.add_paragraph()
        p.text = phase['weeks']
        p.font.size = Pt(20)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        # Goal
        p = tf.add_paragraph()
        p.text = f"\nGoal: {phase['goal']}"
        p.font.size = Pt(22)
        p.font.bold = True

        # Deliverables
        p = tf.add_paragraph()
        p.text = "\nKey Deliverables:"
        p.font.size = Pt(20)
        p.font.bold = True

        for deliverable in phase['deliverables']:
            p = tf.add_paragraph()
            p.text = deliverable
            p.level = 1
            p.font.size = Pt(18)


def add_closing_slide(prs):
    """Add closing slide with next steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Let's Make Piano Keys Production-Ready!"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 169, 224)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Next Steps:"
    p.font.size = Pt(28)
    p.font.bold = True

    next_steps = [
        "Review and approve this plan",
        "Start with Phase 1 (Security)",
        "Create feature branch",
        "Track progress weekly",
        "Launch in 12 weeks! 🎉"
    ]

    for step in next_steps:
        p = tf.add_paragraph()
        p.text = step
        p.level = 1
        p.font.size = Pt(22)


if __name__ == "__main__":
    # Test generation
    source = Path(__file__).parent.parent / ".claude/plans/production-readiness-plan-source.md"
    output = Path(__file__).parent.parent / "output"
    output.mkdir(exist_ok=True)

    result = create_powerpoint_presentation(source, output)
    print(f"✅ PowerPoint presentation generated: {result}")
