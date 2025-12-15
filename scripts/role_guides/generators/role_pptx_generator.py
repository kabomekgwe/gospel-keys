#!/usr/bin/env python3
"""
Role-Specific PowerPoint Presentation Generator
Generates 30-50 slide presentations for each engineering role
"""

from pathlib import Path
from datetime import date
from typing import List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import sys

# Add parent directories to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent / "content"))

from code_snippet_extractor import CodeSnippetExtractor
from role_mapper import RoleMapper


class RolePptxGenerator:
    """Generate role-specific PowerPoint presentations."""

    def __init__(
        self,
        project_root: Path,
        plan_file: Path,
        role_configs_dir: Path,
        output_dir: Path
    ):
        """Initialize the presentation generator."""
        self.extractor = CodeSnippetExtractor(project_root, plan_file, role_configs_dir)
        self.mapper = RoleMapper(plan_file, role_configs_dir)
        self.output_dir = output_dir
        self.output_dir.mkdir(exist_ok=True)

    def generate_role_presentation(self, role_name: str) -> Path:
        """
        Generate PowerPoint presentation for a specific role.

        Args:
            role_name: Name of the role

        Returns:
            Path to generated presentation
        """
        print(f"\n📊 Generating presentation for: {role_name}")
        print("-" * 70)

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Get role summary
        summary = self.mapper.generate_role_summary(role_name)

        # Add slides
        self._add_title_slide(prs, role_name, summary)
        self._add_agenda_slide(prs, summary)
        self._add_overview_slide(prs, summary)

        # Phase slides
        phases = summary.get('phases', [])
        for phase in phases:
            self._add_phase_slide(prs, phase, role_name)

        # Implementation slides
        tasks = summary.get('tasks', [])
        critical_tasks = [t for t in tasks if t.get('priority') == 'CRITICAL'][:5]

        for task in critical_tasks:
            self._add_task_slide(prs, task, role_name)

        # Success metrics slide
        self._add_metrics_slide(prs, summary)

        # Technology stack slide
        self._add_tech_stack_slide(prs, summary)

        # Timeline slide
        self._add_timeline_slide(prs, summary)

        # Closing slide
        self._add_closing_slide(prs, role_name)

        # Save presentation
        filename = f"{role_name.lower().replace(' ', '_').replace('/', '_')}_presentation.pptx"
        output_path = self.output_dir / filename
        prs.save(output_path)

        print(f"✅ Generated: {output_path}")
        print(f"   File size: {output_path.stat().st_size // 1024} KB")
        print(f"   Slides: {len(prs.slides)}")

        return output_path

    def _add_title_slide(self, prs, role_name: str, summary: dict):
        """Add title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = role_name
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 71, 136)
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(8), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Production Readiness Implementation"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(46, 92, 138)
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Date
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(8), Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_frame.text = date.today().strftime('%B %d, %Y')
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(16)
        date_para.alignment = PP_ALIGN.CENTER

    def _add_agenda_slide(self, prs, summary: dict):
        """Add agenda slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = "Agenda"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 71, 136)

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()

        agenda_items = [
            "Role Overview",
            "Implementation Phases",
            "Critical Tasks",
            "Success Metrics",
            "Technology Stack",
            "Timeline & Estimates"
        ]

        for item in agenda_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(24)

    def _add_overview_slide(self, prs, summary: dict):
        """Add role overview slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = "Role Overview"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 71, 136)

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()

        # Description
        desc = summary.get('description', '')
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(18)

        # Key responsibilities
        p = tf.add_paragraph()
        p.text = "\nKey Responsibilities:"
        p.font.size = Pt(22)
        p.font.bold = True

        phases = summary.get('phases', [])
        for phase in phases[:3]:  # Show first 3 phases
            p = tf.add_paragraph()
            p.text = f"Phase {phase['phase']}: {phase['name']}"
            p.level = 1
            p.font.size = Pt(18)

    def _add_phase_slide(self, prs, phase: dict, role_name: str):
        """Add slide for a specific phase."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        phase_num = phase.get('phase')
        phase_name = phase.get('name')
        weeks = phase.get('weeks')

        title = slide.shapes.title
        title.text = f"Phase {phase_num}: {phase_name}"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 71, 136)

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()

        # Weeks
        p = tf.add_paragraph()
        p.text = f"Timeline: Weeks {weeks}"
        p.font.size = Pt(20)
        p.font.italic = True

        # Focus areas
        p = tf.add_paragraph()
        p.text = "\nFocus Areas:"
        p.font.size = Pt(22)
        p.font.bold = True

        focus_areas = phase.get('focus_areas', [])
        for area in focus_areas[:5]:  # Limit to 5 for readability
            p = tf.add_paragraph()
            p.text = area
            p.level = 1
            p.font.size = Pt(16)

    def _add_task_slide(self, prs, task: dict, role_name: str):
        """Add slide for a critical task."""
        snippet = self.extractor.extract_task_code(task)

        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = f"🔴 CRITICAL: {snippet['task']}"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 53, 69)

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()

        # File and priority
        p = tf.add_paragraph()
        p.text = f"File: {snippet['file']}"
        p.font.size = Pt(18)

        if snippet.get('target_lines'):
            p = tf.add_paragraph()
            p.text = f"Lines: {snippet['target_lines']}"
            p.font.size = Pt(18)

        # Task description
        p = tf.add_paragraph()
        p.text = "\nWhat needs to be done:"
        p.font.size = Pt(20)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = snippet['task']
        p.level = 1
        p.font.size = Pt(16)

        # Impact (if it's a security vulnerability)
        if task.get('vulnerability'):
            p = tf.add_paragraph()
            p.text = f"\n⚠️ Vulnerability: {task['vulnerability']}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(220, 53, 69)

    def _add_metrics_slide(self, prs, summary: dict):
        """Add success metrics slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = "Success Metrics"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 169, 224)

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()

        metrics = summary.get('success_metrics', [])

        for metric in metrics:
            p = tf.add_paragraph()
            p.text = f"✓ {metric['metric']}: {metric['target']}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(40, 167, 69)

    def _add_tech_stack_slide(self, prs, summary: dict):
        """Add technology stack slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = "Technology Stack"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 71, 136)

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()

        tech_stack = summary.get('tech_stack', [])

        for tech in tech_stack[:12]:  # Limit to 12 for readability
            p = tf.add_paragraph()
            p.text = tech
            p.font.size = Pt(18)

    def _add_timeline_slide(self, prs, summary: dict):
        """Add timeline slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = "Timeline & Estimates"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 71, 136)

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()

        time_estimates = summary.get('time_estimates', {})

        for phase, estimate in time_estimates.items():
            p = tf.add_paragraph()
            p.text = f"{phase}: {estimate}"
            p.font.size = Pt(20)

    def _add_closing_slide(self, prs, role_name: str):
        """Add closing slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        title.text = "Let's Build Production-Ready Software!"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 169, 224)

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()

        p = tf.add_paragraph()
        p.text = "Next Steps:"
        p.font.size = Pt(28)
        p.font.bold = True

        next_steps = [
            "Review this implementation guide",
            "Set up your development environment",
            "Start with Phase 1 critical tasks",
            "Track progress with success metrics",
            "Collaborate with other roles"
        ]

        for step in next_steps:
            p = tf.add_paragraph()
            p.text = step
            p.level = 1
            p.font.size = Pt(20)

    def generate_all_roles(self) -> List[Path]:
        """Generate presentations for all roles."""
        generated_ppts = []

        roles = self.mapper.get_all_roles()

        print(f"\n🚀 Generating presentations for {len(roles)} roles...")
        print("=" * 70)

        for role in roles:
            ppt_path = self.generate_role_presentation(role)
            generated_ppts.append(ppt_path)

        return generated_ppts


if __name__ == "__main__":
    # Test presentation generation
    from pathlib import Path

    project_root = Path(__file__).parent.parent.parent.parent
    plan_file = project_root / ".claude/plans/production-readiness-plan-source.md"
    configs_dir = Path(__file__).parent.parent / "templates/role_configs"
    output_dir = project_root / "output/role_guides"

    generator = RolePptxGenerator(project_root, plan_file, configs_dir, output_dir)

    # Generate Backend Engineer presentation as test
    print("Testing Role Presentation Generator...")
    print("=" * 70)

    ppt_path = generator.generate_role_presentation("Backend Engineer")

    print(f"\n✅ Test complete! Presentation generated at:")
    print(f"   {ppt_path}")
